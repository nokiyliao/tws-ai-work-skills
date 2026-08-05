#!/usr/bin/env python3
"""Cross-platform isolated runtime bootstrap for TWS AI deck production.

Detects Windows/macOS, Python, uv, an isolated runtime under
~/.codex/runtimes/tws-ai (override with TWS_AI_RUNTIME_HOME), installs pinned
Python packages, and fail-closes when system renderer/OCR tools are missing.
Never installs Office, LibreOffice, or Tesseract; those are reported as typed
blockers when absent.

Modes:
  install  create/update isolated runtime, install lockfile, smoke, remote preflight
  check    verify readiness without installing packages
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import platform
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any, Callable


SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_ROOT = SCRIPT_DIR.parent
RUNTIME_SPEC_DIR = SKILL_ROOT / "runtime"
REQUIREMENTS_LOCK = RUNTIME_SPEC_DIR / "requirements.lock"
REQUIREMENTS_TXT = RUNTIME_SPEC_DIR / "requirements.txt"
DEFAULT_RUNTIME_HOME = Path.home() / ".codex" / "runtimes" / "tws-ai"
MIN_PYTHON = (3, 10)

BLOCKER_PYTHON = "PYTHON_VERSION"
BLOCKER_UV = "UV_MISSING"
BLOCKER_VENV = "RUNTIME_VENV"
BLOCKER_IMPORTS = "PYTHON_IMPORTS"
BLOCKER_RENDERER = "SYSTEM_RENDERER_MISSING"
BLOCKER_RASTERIZER = "PDF_RASTERIZER_MISSING"
BLOCKER_OCR = "SYSTEM_OCR_MISSING"
BLOCKER_SMOKE = "SMOKE_TEST"
BLOCKER_REMOTE = "REMOTE_ASSET_BOOTSTRAP"
BLOCKER_PREFLIGHT = "PREFLIGHT"
BLOCKER_UNSUPPORTED = "UNSUPPORTED_PLATFORM"
BLOCKER_RUNTIME_STATE = "RUNTIME_UNVERIFIED"

# Aliases used by some call sites / tests.
BLOCKER_PYTHON_VERSION = BLOCKER_PYTHON
BLOCKER_UV_MISSING = BLOCKER_UV
BLOCKER_RUNTIME_MISSING = BLOCKER_VENV
BLOCKER_IMPORT_MISSING = BLOCKER_IMPORTS
BLOCKER_RENDERER_MISSING = BLOCKER_RENDERER
BLOCKER_RASTERIZER_MISSING = BLOCKER_RASTERIZER
BLOCKER_OCR_MISSING = BLOCKER_OCR
BLOCKER_UNSUPPORTED_OS = BLOCKER_UNSUPPORTED


def runtime_home(override: Path | None = None) -> Path:
    if override is not None:
        return override.expanduser().resolve()
    env = os.environ.get("TWS_AI_RUNTIME_HOME", "").strip()
    if env:
        return Path(env).expanduser().resolve()
    return DEFAULT_RUNTIME_HOME


def runtime_paths(home: Path | None = None) -> dict[str, Path]:
    root = Path(home) if home is not None else runtime_home()
    venv_dir = root / "venv"
    if sys.platform == "win32":
        python = venv_dir / "Scripts" / "python.exe"
        pip = venv_dir / "Scripts" / "pip.exe"
    else:
        python = venv_dir / "bin" / "python"
        pip = venv_dir / "bin" / "pip"
    return {
        "home": root,
        "root": root,
        "venv": venv_dir,
        "python": python,
        "pip": pip,
        "marker": root / "runtime-manifest.json",
        "uv_cache": root / "uv-cache",
    }


def detect_platform(system: str | None = None) -> dict[str, Any]:
    raw = system if system is not None else platform.system()
    lowered = raw.lower()
    if lowered in {"darwin", "macos"}:
        family = "macos"
    elif lowered in {"windows", "win32"}:
        family = "windows"
    elif lowered == "linux":
        family = "linux"
    else:
        family = lowered or "unknown"
    # Learners are macOS/Windows; Linux is accepted for CI/tooling probes.
    ok = family in {"macos", "windows", "linux"}
    return {
        "ok": ok,
        "family": family,
        "os": family,
        "system": raw,
        "machine": platform.machine(),
        "release": platform.release(),
        "platform": sys.platform,
        "blocker": None if ok else BLOCKER_UNSUPPORTED,
        "detail": None if ok else f"unsupported platform: {family}",
    }


def detect_python(
    version_info: tuple[int, ...] | None = None,
    min_version: tuple[int, int] = MIN_PYTHON,
    executable: str | None = None,
) -> dict[str, Any]:
    exe = executable or sys.executable
    if version_info is not None:
        version = (int(version_info[0]), int(version_info[1]), int(version_info[2]) if len(version_info) > 2 else 0)
    elif executable is None or executable == sys.executable:
        version = sys.version_info[:3]
    else:
        try:
            out = subprocess.check_output(
                [exe, "-c", "import sys; print('.'.join(map(str, sys.version_info[:3])))"],
                text=True,
                timeout=30,
            ).strip()
            parts = tuple(int(p) for p in out.split("."))
            version = (parts + (0, 0, 0))[:3]
        except (OSError, subprocess.SubprocessError, ValueError) as exc:
            return {
                "ok": False,
                "executable": exe,
                "version": None,
                "blocker": BLOCKER_PYTHON,
                "detail": f"cannot probe python: {exc}",
            }
    ok = version[:2] >= min_version
    return {
        "ok": ok,
        "executable": exe,
        "version": f"{version[0]}.{version[1]}.{version[2]}",
        "min_version": f"{min_version[0]}.{min_version[1]}",
        "blocker": None if ok else BLOCKER_PYTHON,
        "detail": None if ok else f"python {version[0]}.{version[1]} < {min_version[0]}.{min_version[1]}",
    }


def which_tool(name: str, path_env: str | None = None) -> str | None:
    if path_env is None:
        return shutil.which(name)
    return shutil.which(name, path=path_env)


def which(name: str, path: str | None = None) -> str | None:
    return which_tool(name, path_env=path)


def detect_uv(
    path_env: str | None = None,
    which_fn: Callable[..., str | None] | None = None,
) -> dict[str, Any]:
    if which_fn is not None:
        path = which_fn("uv", path_env)
    else:
        path = which_tool("uv", path_env=path_env)
    if not path:
        candidates = [
            Path.home() / ".local" / "bin" / ("uv.exe" if sys.platform == "win32" else "uv"),
            Path.home() / ".cargo" / "bin" / ("uv.exe" if sys.platform == "win32" else "uv"),
        ]
        for candidate in candidates:
            if candidate.is_file():
                path = str(candidate)
                break
    if not path:
        return {
            "ok": False,
            "path": None,
            "version": None,
            "blocker": BLOCKER_UV,
            "detail": "uv not found; Codex may install uv into ~/.local/bin without admin rights",
        }
    version = None
    try:
        version = subprocess.check_output([path, "--version"], text=True, timeout=30).strip()
    except (OSError, subprocess.SubprocessError):
        version = "unknown"
    return {"ok": True, "path": path, "version": version, "blocker": None, "detail": None}


def detect_renderer(
    *,
    which: Callable[[str], str | None] | None = None,
    which_fn: Callable[..., str | None] | None = None,
    platform_name: str | None = None,
    win32com_available: bool | None = None,
    powerpoint_probe: Callable[[], bool] | None = None,
) -> dict[str, Any]:
    plat = platform_name or sys.platform
    engines: list[str] = []

    def find(name: str) -> str | None:
        if which_fn is not None:
            return which_fn(name, None)
        if which is not None:
            return which(name)
        return shutil.which(name)

    if plat in {"win32", "windows"}:
        if powerpoint_probe is not None:
            win32com_available = powerpoint_probe()
        elif win32com_available is None:
            try:
                import win32com.client  # type: ignore  # noqa: F401

                win32com_available = True
            except Exception:
                win32com_available = False
        if win32com_available:
            engines.append("powerpoint-com")

    for name in ("soffice", "libreoffice"):
        path = find(name)
        if path:
            engines.append(f"libreoffice:{path}")
            break

    if engines:
        return {
            "ok": True,
            "engines": engines,
            "primary": engines[0],
            "blocker": None,
            "detail": None,
            "admin_or_gui": False,
            "admin_or_gui_required": False,
        }

    detail = (
        "no PPTX renderer available; install Microsoft PowerPoint (Windows COM) "
        "or LibreOffice (soffice). This tool is not auto-installed by bootstrap."
    )
    return {
        "ok": False,
        "engines": [],
        "primary": None,
        "blocker": BLOCKER_RENDERER,
        "detail": detail,
        "admin_or_gui": True,
        "admin_or_gui_required": True,
    }


def detect_rasterizer(
    python_exe: Path | None = None,
    *,
    which: Callable[[str], str | None] | None = None,
    which_fn: Callable[..., str | None] | None = None,
    fitz_available: bool | None = None,
    import_probe: Callable[[str, Path | None], bool] | None = None,
) -> dict[str, Any]:
    engines: list[str] = []

    def find(name: str) -> str | None:
        if which_fn is not None:
            return which_fn(name, None)
        if which is not None:
            return which(name)
        return shutil.which(name)

    pdftoppm = find("pdftoppm")
    if pdftoppm:
        engines.append(f"pdftoppm:{pdftoppm}")

    if fitz_available is None:
        if import_probe is not None:
            fitz_available = import_probe("fitz", python_exe)
        else:
            fitz_available = _module_importable("fitz", python_exe)
    if fitz_available:
        engines.append("pymupdf")

    if engines:
        primary = engines[0] if engines[0] != "pymupdf" else "pymupdf"
        if engines[0] == "pymupdf" or (len(engines) == 1 and engines[0] == "pymupdf"):
            primary = "pymupdf"
        elif "pymupdf" in engines and not pdftoppm:
            primary = "pymupdf"
        else:
            primary = engines[0]
        # Prefer short name for pymupdf-only case.
        if engines == ["pymupdf"]:
            primary = "pymupdf"
        return {
            "ok": True,
            "engines": engines,
            "primary": primary if primary != "pymupdf" and not primary.startswith("pdftoppm") else (
                "pymupdf" if "pymupdf" in engines and primary == "pymupdf" else engines[0]
            ),
            "blocker": None,
            "detail": None,
        }

    return {
        "ok": False,
        "engines": [],
        "primary": None,
        "blocker": BLOCKER_RASTERIZER,
        "detail": "install Poppler pdftoppm or ensure PyMuPDF is installed in the isolated runtime",
    }


def detect_ocr(
    *,
    which: Callable[[str], str | None] | None = None,
    which_fn: Callable[..., str | None] | None = None,
    platform_name: str | None = None,
) -> dict[str, Any]:
    plat = platform_name or sys.platform

    def find(name: str) -> str | None:
        if which_fn is not None:
            return which_fn(name, None)
        if which is not None:
            return which(name)
        return shutil.which(name)

    engines: list[str] = []
    if plat in {"darwin", "macos"} and find("swiftc"):
        engines.append("macos-vision")
    tesseract = find("tesseract")
    if tesseract:
        engines.append(f"tesseract:{tesseract}")

    if engines:
        return {
            "ok": True,
            "engines": engines,
            "primary": engines[0],
            "blocker": None,
            "detail": None,
            "admin_or_gui": False,
            "admin_or_gui_required": False,
        }

    detail = (
        "no OCR engine available; on macOS install Xcode CLT (swiftc) for Vision, "
        "or install Tesseract on Windows/macOS. OCR engines are not auto-installed."
    )
    return {
        "ok": False,
        "engines": [],
        "primary": None,
        "blocker": BLOCKER_OCR,
        "detail": detail,
        "admin_or_gui": True,
        "admin_or_gui_required": True,
    }


def _module_importable(module: str, python_exe: Path | None = None) -> bool:
    if python_exe is None or Path(python_exe) == Path(sys.executable):
        try:
            __import__(module)
            return True
        except Exception:
            return False
    if not Path(python_exe).is_file():
        return False
    try:
        subprocess.check_call(
            [str(python_exe), "-c", f"import {module}"],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=60,
        )
        return True
    except (OSError, subprocess.SubprocessError):
        return False


def detect_imports(
    python_exe: Path | None = None,
    platform_name: str | None = None,
    import_probe: Callable[[str, Path | None], bool] | None = None,
) -> dict[str, Any]:
    required = {
        "pptx": "python-pptx",
        "PIL": "Pillow",
        "fitz": "PyMuPDF",
    }
    plat = platform_name or sys.platform
    if plat in {"win32", "windows"}:
        required["win32com"] = "pywin32"
    probe = import_probe or (lambda mod, exe: _module_importable(mod, exe))
    modules = {mod: probe(mod, python_exe) for mod in required}
    missing = [pkg for mod, pkg in required.items() if not modules[mod]]
    ok = not missing
    return {
        "ok": ok,
        "modules": modules,
        "missing": missing,
        "python": str(python_exe) if python_exe else sys.executable,
        "blocker": None if ok else BLOCKER_IMPORTS,
        "detail": None if ok else f"missing imports: {', '.join(missing)}",
    }


def detect_runtime(home: Path | None = None) -> dict[str, Any]:
    paths = runtime_paths(home)
    python_ok = paths["python"].is_file()
    manifest: dict[str, Any] = {}
    if paths["marker"].is_file():
        try:
            manifest = json.loads(paths["marker"].read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            manifest = {}
    lock_digest = hashlib.sha256(REQUIREMENTS_LOCK.read_bytes()).hexdigest() if REQUIREMENTS_LOCK.is_file() else None
    manifest_ok = (
        manifest.get("schema_version") == "tws_ai_runtime_manifest_v1"
        and manifest.get("status") == "PASS"
        and manifest.get("requirements_sha256") == lock_digest
    )
    ok = python_ok and manifest_ok
    return {
        "ok": ok,
        "home": str(paths["home"]),
        "python": str(paths["python"]) if python_ok else None,
        "venv": str(paths["venv"]),
        "manifest": str(paths["marker"]),
        "blocker": None if ok else (BLOCKER_VENV if not python_ok else BLOCKER_RUNTIME_STATE),
        "detail": None if ok else (
            f"missing isolated runtime python: {paths['python']}"
            if not python_ok
            else "runtime manifest is missing, failed, or does not match the current dependency lock"
        ),
    }


def ensure_venv(paths: dict[str, Path], *, uv_path: str | None = None) -> dict[str, Any]:
    home = paths["home"]
    venv_dir = paths["venv"]
    home.mkdir(parents=True, exist_ok=True)
    if paths["python"].is_file():
        return {"ok": True, "created": False, "venv": str(venv_dir), "blocker": None, "detail": None}
    if not uv_path:
        return {
            "ok": False,
            "created": False,
            "venv": str(venv_dir),
            "blocker": BLOCKER_UV,
            "detail": "uv is required to create the managed isolated runtime",
        }
    try:
        subprocess.check_call(
            [uv_path, "venv", str(venv_dir), "--python", f"{MIN_PYTHON[0]}.{MIN_PYTHON[1]}"],
            timeout=180,
        )
    except (OSError, subprocess.SubprocessError) as exc:
        return {
            "ok": False,
            "created": False,
            "venv": str(venv_dir),
            "blocker": BLOCKER_VENV,
            "detail": f"failed to create isolated venv: {exc}",
        }
    if not paths["python"].is_file():
        return {
            "ok": False,
            "created": False,
            "venv": str(venv_dir),
            "blocker": BLOCKER_VENV,
            "detail": f"venv python missing after create: {paths['python']}",
        }
    return {"ok": True, "created": True, "venv": str(venv_dir), "blocker": None, "detail": None}


def install_packages(paths: dict[str, Path], *, uv_path: str | None = None) -> dict[str, Any]:
    req = REQUIREMENTS_LOCK if REQUIREMENTS_LOCK.is_file() else REQUIREMENTS_TXT
    if not req.is_file():
        return {
            "ok": False,
            "blocker": BLOCKER_IMPORTS,
            "detail": f"requirements file missing: {req}",
            "requirements": str(req),
        }
    python = paths["python"]
    if not uv_path:
        return {
            "ok": False,
            "blocker": BLOCKER_UV,
            "detail": "uv is required to install the managed runtime",
            "requirements": str(req),
        }
    try:
        subprocess.check_call(
            [uv_path, "pip", "install", "--python", str(python), "-r", str(req)],
            timeout=600,
        )
    except (OSError, subprocess.SubprocessError) as exc:
        return {
            "ok": False,
            "blocker": BLOCKER_IMPORTS,
            "detail": f"package install failed: {exc}",
            "requirements": str(req),
        }
    return {
        "ok": True,
        "blocker": None,
        "detail": None,
        "requirements": str(req),
        "python": str(python),
    }


def write_manifest(paths: dict[str, Path], report: dict[str, Any]) -> None:
    paths["home"].mkdir(parents=True, exist_ok=True)
    payload = {
        "schema_version": "tws_ai_runtime_manifest_v1",
        "runtime_home": str(paths["home"]),
        "python": str(paths["python"]),
        "status": report.get("status"),
        "requirements_sha256": hashlib.sha256(REQUIREMENTS_LOCK.read_bytes()).hexdigest(),
        "platform": report.get("checks", {}).get("platform"),
        "imports": report.get("checks", {}).get("imports"),
    }
    paths["marker"].write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")


def smoke_pptx(python_exe: Path, workdir: Path) -> dict[str, Any]:
    pptx_path = workdir / "smoke.pptx"
    script = (
        "from pptx import Presentation\n"
        "from pptx.util import Inches, Pt\n"
        f"out = {str(pptx_path)!r}\n"
        "prs = Presentation()\n"
        "prs.slide_width = Inches(13.333)\n"
        "prs.slide_height = Inches(7.5)\n"
        "slide = prs.slides.add_slide(prs.slide_layouts[6])\n"
        "box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(10), Inches(1.2))\n"
        "tf = box.text_frame\n"
        "tf.text = 'TWS AI runtime smoke'\n"
        "tf.paragraphs[0].font.size = Pt(28)\n"
        "prs.save(out)\n"
        "print(out)\n"
    )
    try:
        subprocess.check_call([str(python_exe), "-c", script], timeout=120, cwd=str(workdir))
        ok = pptx_path.is_file() and pptx_path.stat().st_size > 1000
        return {
            "ok": ok,
            "path": str(pptx_path),
            "blocker": None if ok else BLOCKER_SMOKE,
            "detail": None if ok else "pptx file was not created",
        }
    except (OSError, subprocess.SubprocessError) as exc:
        return {"ok": False, "path": str(pptx_path), "blocker": BLOCKER_SMOKE, "detail": str(exc)}


def smoke_render(python_exe: Path, workdir: Path, pptx_path: Path) -> dict[str, Any]:
    renderer = SKILL_ROOT.parent / "nokiy-presentation-generator" / "scripts" / "render_slides.py"
    output = workdir / "rendered"
    if not renderer.is_file():
        return {"ok": False, "path": None, "blocker": BLOCKER_SMOKE, "detail": f"renderer script missing: {renderer}"}
    try:
        result = subprocess.run(
            [str(python_exe), str(renderer), str(pptx_path), "--output", str(output)],
            text=True,
            capture_output=True,
            timeout=180,
        )
        images = sorted(output.glob("slide-*.png"))
        ok = result.returncode == 0 and bool(images) and all(path.stat().st_size > 100 for path in images)
        return {
            "ok": ok,
            "path": str(images[0]) if images else None,
            "blocker": None if ok else BLOCKER_SMOKE,
            "detail": None if ok else ((result.stderr or result.stdout or "PPTX render produced no PNG").strip()),
        }
    except (OSError, subprocess.SubprocessError) as exc:
        return {"ok": False, "path": None, "blocker": BLOCKER_SMOKE, "detail": str(exc)}


def smoke_ocr(image_path: Path, ocr_info: dict[str, Any] | Any) -> dict[str, Any]:
    # Accept dict or CheckResult-like objects.
    if hasattr(ocr_info, "ok"):
        ok_flag = bool(ocr_info.ok)
        primary = str((getattr(ocr_info, "data", None) or {}).get("primary") or "")
        detail = getattr(ocr_info, "detail", None) or "OCR engine missing"
        blocker = getattr(ocr_info, "blocker", None) or BLOCKER_OCR
        engines = (getattr(ocr_info, "data", None) or {}).get("engines") or []
    else:
        ok_flag = bool(ocr_info.get("ok"))
        primary = str(ocr_info.get("primary") or "")
        detail = ocr_info.get("detail") or "OCR engine missing"
        blocker = ocr_info.get("blocker") or BLOCKER_OCR
        engines = ocr_info.get("engines") or []

    if not ok_flag:
        return {
            "ok": False,
            "blocker": blocker,
            "detail": detail,
            "admin_or_gui": True,
            "engine": None,
        }

    # Normalize primary from engines list.
    if not primary and engines:
        primary = str(engines[0])
    try:
        if primary.startswith("tesseract:") or primary == "tesseract":
            binary = primary.split(":", 1)[1] if ":" in primary else which_tool("tesseract")
            if not binary:
                return {"ok": False, "blocker": BLOCKER_SMOKE, "detail": "tesseract binary missing", "engine": "tesseract"}
            res = subprocess.run(
                [binary, str(image_path), "stdout", "-l", "eng"],
                text=True,
                capture_output=True,
                timeout=120,
            )
            text = res.stdout or ""
            ok = res.returncode == 0  # empty text still proves engine runs for fixture images
            return {
                "ok": ok,
                "engine": "tesseract",
                "text_preview": text.strip()[:120],
                "blocker": None if ok else BLOCKER_SMOKE,
                "detail": None if ok else (res.stderr or "tesseract failed"),
            }
        if primary in {"macos-vision", "vision_swift"} or str(primary).startswith("vision_swift"):
            ocr_script = SKILL_ROOT.parent / "nokiy-presentation-generator" / "scripts" / "ocr_rendered_slides.py"
            result = subprocess.run(
                [sys.executable, str(ocr_script), str(image_path), "--engine", "vision", "--json"],
                text=True,
                capture_output=True,
                timeout=180,
            )
            try:
                rows = json.loads(result.stdout or "[]")
            except json.JSONDecodeError:
                rows = []
            ok = result.returncode == 0 and len(rows) == 1 and not rows[0].get("error")
            return {
                "ok": ok,
                "engine": "macos-vision",
                "text_preview": str(rows[0].get("text", ""))[:120] if rows else None,
                "blocker": None if ok else BLOCKER_SMOKE,
                "detail": None if ok else ((result.stderr or result.stdout or "Vision OCR smoke failed").strip()),
            }
        return {
            "ok": False,
            "engine": primary,
            "blocker": BLOCKER_OCR,
            "detail": f"unsupported OCR engine: {primary}",
        }
    except (OSError, subprocess.SubprocessError) as exc:
        return {"ok": False, "blocker": BLOCKER_SMOKE, "detail": str(exc)}


def run_remote_bootstrap(config_path: Path | None = None, *, skip: bool = False) -> dict[str, Any]:
    if skip:
        return {"ok": True, "skipped": True, "blocker": None, "detail": "skipped"}
    bootstrap = SCRIPT_DIR / "bootstrap_learner.py"
    if not bootstrap.is_file():
        return {"ok": False, "blocker": BLOCKER_REMOTE, "detail": f"missing {bootstrap}"}
    cmd = [sys.executable, str(bootstrap)]
    if config_path is not None:
        cmd.extend(["--config", str(config_path)])
    try:
        res = subprocess.run(cmd, text=True, capture_output=True, timeout=90)
    except (OSError, subprocess.SubprocessError) as exc:
        return {"ok": False, "blocker": BLOCKER_REMOTE, "detail": str(exc)}
    ok = res.returncode == 0
    return {
        "ok": ok,
        "returncode": res.returncode,
        "stdout": (res.stdout or "").strip()[:500],
        "stderr": (res.stderr or "").strip()[:500],
        "blocker": None if ok else BLOCKER_REMOTE,
        "detail": None if ok else ((res.stderr or res.stdout or "remote bootstrap failed").strip()),
    }


def run_preflight_script(workflow: str = "tws-new-factory", *, skip: bool = False) -> dict[str, Any]:
    if skip:
        return {"ok": True, "skipped": True, "blocker": None, "detail": "skipped"}
    preflight = SCRIPT_DIR / "preflight.py"
    if not preflight.is_file():
        return {"ok": False, "blocker": BLOCKER_PREFLIGHT, "detail": f"missing {preflight}"}
    try:
        res = subprocess.run(
            [sys.executable, str(preflight), "--workflow", workflow, "--asset-mode", "remote"],
            text=True,
            capture_output=True,
            timeout=90,
        )
    except (OSError, subprocess.SubprocessError) as exc:
        return {"ok": False, "blocker": BLOCKER_PREFLIGHT, "detail": str(exc)}
    payload: dict[str, Any] = {}
    try:
        payload = json.loads(res.stdout or "{}")
    except json.JSONDecodeError:
        payload = {"raw_stdout": (res.stdout or "")[:500]}
    ok = res.returncode == 0 and payload.get("status") == "PASS"
    return {
        "ok": ok,
        "returncode": res.returncode,
        "report": payload,
        "blocker": None if ok else BLOCKER_PREFLIGHT,
        "detail": None if ok else ((res.stderr or "preflight failed").strip() or str(payload.get("failed"))),
    }


def collect_blockers(checks: dict[str, Any]) -> list[dict[str, Any]]:
    blockers: list[dict[str, Any]] = []
    for name, result in checks.items():
        if not isinstance(result, dict):
            continue
        if result.get("ok") is False and result.get("blocker"):
            blockers.append(
                {
                    "check": name,
                    "blocker": result.get("blocker"),
                    "detail": result.get("detail"),
                    "admin_or_gui": bool(result.get("admin_or_gui") or result.get("admin_or_gui_required")),
                }
            )
    return blockers


def build_report(
    *,
    mode: str,
    home: Path,
    checks: dict[str, Any],
) -> dict[str, Any]:
    blockers = collect_blockers(checks)
    paths = runtime_paths(home)
    return {
        "status": "PASS" if not blockers else "FAIL",
        "mode": mode,
        "runtime_home": str(paths["home"]),
        "runtime_python": str(paths["python"]) if paths["python"].is_file() else None,
        "checks": checks,
        "blockers": blockers,
        "failed": [b["check"] for b in blockers],
        "notes": [
            "Python packages install only into the isolated runtime; the global interpreter is never modified.",
            "Office, LibreOffice, and Tesseract are never auto-installed; missing tools are typed blockers.",
            "Student assets remain remote-only via the published HTTPS Tunnel API.",
        ],
    }


# Keep notes free of the substring "global-packages-avoided" so reports cannot be mistaken
# for a global pip install path.


def run_check(home: Path, *, skip_remote: bool = False) -> dict[str, Any]:
    paths = runtime_paths(home)
    python_for_imports = paths["python"] if paths["python"].is_file() else None
    imports = detect_imports(python_for_imports)
    checks: dict[str, Any] = {
        "platform": detect_platform(),
        "python_host": detect_python(),
        "uv": detect_uv(),
        "runtime": detect_runtime(home),
        "venv": {
            "ok": paths["python"].is_file(),
            "path": str(paths["venv"]),
            "python": str(paths["python"]) if paths["python"].is_file() else None,
            "blocker": None if paths["python"].is_file() else BLOCKER_VENV,
            "detail": None if paths["python"].is_file() else f"missing isolated runtime python: {paths['python']}",
        },
        "imports": imports,
        "renderer": detect_renderer(win32com_available=imports["modules"].get("win32com")),
        "rasterizer": detect_rasterizer(python_exe=python_for_imports),
        "ocr": detect_ocr(),
    }
    if not skip_remote:
        checks["remote_bootstrap"] = run_remote_bootstrap()
        checks["preflight"] = run_preflight_script()
    return build_report(mode="check", home=home, checks=checks)


def run_install(home: Path, *, skip_remote: bool = False, skip_smoke: bool = False) -> dict[str, Any]:
    paths = runtime_paths(home)
    checks: dict[str, Any] = {
        "platform": detect_platform(),
        "python_host": detect_python(),
        "uv": detect_uv(),
    }
    if not checks["platform"]["ok"] or not checks["python_host"]["ok"]:
        return build_report(mode="install", home=home, checks=checks)

    if not checks["uv"]["ok"]:
        return build_report(mode="install", home=home, checks=checks)
    uv_path = checks["uv"]["path"]
    checks["venv"] = ensure_venv(paths, uv_path=uv_path)
    if not checks["venv"]["ok"]:
        return build_report(mode="install", home=home, checks=checks)

    checks["packages"] = install_packages(paths, uv_path=uv_path)
    if not checks["packages"]["ok"]:
        return build_report(mode="install", home=home, checks=checks)

    checks["imports"] = detect_imports(paths["python"])
    checks["renderer"] = detect_renderer(win32com_available=checks["imports"]["modules"].get("win32com"))
    checks["rasterizer"] = detect_rasterizer(python_exe=paths["python"])
    checks["ocr"] = detect_ocr()

    if not skip_smoke:
        with tempfile.TemporaryDirectory(prefix="tws-ai-smoke-") as temp:
            workdir = Path(temp)
            checks["smoke_pptx"] = smoke_pptx(paths["python"], workdir)
            pptx = Path(checks["smoke_pptx"].get("path") or workdir / "missing.pptx")
            checks["smoke_render"] = (
                smoke_render(paths["python"], workdir, pptx)
                if checks["smoke_pptx"].get("ok") and pptx.is_file()
                else {"ok": False, "blocker": BLOCKER_SMOKE, "detail": "PPTX smoke did not produce a deck"}
            )
            image = Path(checks["smoke_render"].get("path") or workdir / "missing.png")
            if checks["smoke_render"].get("ok") and image.is_file():
                checks["smoke_ocr"] = smoke_ocr(image, checks["ocr"])
            else:
                checks["smoke_ocr"] = {
                    "ok": False,
                    "blocker": BLOCKER_SMOKE,
                    "detail": "rasterize smoke did not produce an image for OCR",
                }

    # Preflight validates the sealed runtime marker. Write a provisional marker
    # only after local imports/render/OCR checks, then overwrite it with the
    # final remote/preflight result below.
    write_manifest(paths, build_report(mode="install", home=home, checks=checks))
    checks["remote_bootstrap"] = run_remote_bootstrap(skip=skip_remote)
    checks["preflight"] = run_preflight_script(skip=skip_remote)

    report = build_report(mode="install", home=home, checks=checks)
    write_manifest(paths, report)
    return report


def run_bootstrap(
    mode: str,
    runtime_home: Path | None = None,
    *,
    skip_network: bool = False,
    skip_remote: bool | None = None,
    skip_install: bool = False,
    skip_smoke: bool = False,
    **_injected: Any,
) -> dict[str, Any]:
    """Compatibility wrapper. Prefer run_check / run_install."""
    home = Path(runtime_home) if runtime_home is not None else DEFAULT_RUNTIME_HOME
    remote_skip = skip_network if skip_remote is None else skip_remote
    if mode == "install" and not skip_install:
        return run_install(home, skip_remote=remote_skip, skip_smoke=skip_smoke)
    return run_check(home, skip_remote=remote_skip)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "mode",
        nargs="?",
        choices=("install", "check"),
        default="check",
        help="install creates the isolated runtime; check only verifies",
    )
    parser.add_argument(
        "--runtime-home",
        type=Path,
        default=None,
        help="override isolated runtime root (default: ~/.codex/runtimes/tws-ai)",
    )
    parser.add_argument(
        "--skip-remote",
        "--skip-network",
        dest="skip_remote",
        action="store_true",
        help="skip remote catalog bootstrap and preflight (tests / offline)",
    )
    parser.add_argument(
        "--skip-smoke",
        action="store_true",
        help="skip PPTX/PNG/OCR smoke tests",
    )
    parser.add_argument("--json", action="store_true", default=True, help="emit machine-readable JSON")
    args = parser.parse_args(argv)

    home = runtime_home(args.runtime_home)
    if args.mode == "install":
        report = run_install(home, skip_remote=args.skip_remote, skip_smoke=args.skip_smoke)
    else:
        report = run_check(home, skip_remote=args.skip_remote)

    print(json.dumps(report, indent=2, ensure_ascii=True))
    if report["status"] != "PASS":
        admin = [b for b in report["blockers"] if b.get("admin_or_gui")]
        if admin:
            print(
                "FAIL_CLOSED: one or more system tools require admin/GUI install "
                f"({', '.join(b['blocker'] for b in admin)})",
                file=sys.stderr,
            )
        else:
            print(f"FAIL_CLOSED: {', '.join(report['failed'])}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
