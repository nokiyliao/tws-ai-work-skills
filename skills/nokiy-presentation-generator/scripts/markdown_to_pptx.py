#!/usr/bin/env python3
"""Convert constrained Markdown into an editable, business-ready PPTX deck."""

from __future__ import annotations

import argparse
import math
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Helvetica Neue"
FONT_CJK = "PingFang TC"
MIN_BODY = 10.5
MIN_CAPTION = 8.5

THEMES = {
    "industrial": {
        "bg": "F7F8F4",
        "ink": "17352F",
        "deep": "0F2924",
        "text": "263734",
        "muted": "6D7A76",
        "accent": "D96C2C",
        "accent2": "4A7C59",
        "panel": "FFFFFF",
        "panel2": "EEF2EC",
        "rule": "B9C3BC",
        "soft": "E7ECE6",
    },
    "clean": {
        "bg": "FFFFFF",
        "ink": "1F2937",
        "deep": "111827",
        "text": "374151",
        "muted": "6B7280",
        "accent": "2563EB",
        "accent2": "0891B2",
        "panel": "F8FAFC",
        "panel2": "EEF2FF",
        "rule": "CBD5E1",
        "soft": "E5E7EB",
    },
}


@dataclass
class SlideSpec:
    title: str
    body: list[str]
    bullets: list[str]
    images: list[Path]


def split_slides(markdown: str) -> list[str]:
    chunks = re.split(r"(?m)^\s*---\s*$", markdown)
    return [chunk.strip() for chunk in chunks if chunk.strip()]


def parse_slide(chunk: str, base_dir: Path) -> SlideSpec:
    title = ""
    body: list[str] = []
    bullets: list[str] = []
    images: list[Path] = []

    in_comment = False
    for raw in chunk.splitlines():
        line = raw.strip()
        if not line:
            continue
        if "<!--" in line:
            in_comment = True
        if in_comment:
            if "-->" in line:
                in_comment = False
            continue
        image = re.match(r"!\[[^\]]*\]\(([^)]+)\)", line)
        if image:
            path = Path(image.group(1)).expanduser()
            images.append(path if path.is_absolute() else base_dir / path)
            continue
        heading = re.match(r"^#{1,3}\s+(.+)$", line)
        if heading and not title:
            title = heading.group(1).strip()
            continue
        bullet = re.match(r"^(?:[-*]|\d+[.)])\s+(.+)$", line)
        if bullet:
            bullets.append(bullet.group(1).strip())
        else:
            body.append(line)

    if not title:
        title = body.pop(0) if body else "Untitled slide"
    return SlideSpec(title=title, body=body, bullets=bullets, images=images)


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.replace("#", "").upper()
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def visual_units(text: str) -> float:
    units = 0.0
    for ch in str(text):
        if "\u4e00" <= ch <= "\u9fff":
            units += 1.0
        elif ch.isspace():
            units += 0.33
        elif ord(ch) < 128:
            units += 0.55
        else:
            units += 0.80
    return units


def estimate_wrapped_lines(text: str, width_in: float, font_pt: float) -> int:
    capacity = max(4.0, width_in * 72.0 / max(1.0, font_pt) * 0.92)
    return sum(max(1, math.ceil(visual_units(line) / capacity)) for line in str(text).splitlines() or [""])


def safe_font_size(font_size: float) -> float:
    floor = MIN_CAPTION if font_size < 10 else MIN_BODY
    return max(font_size, floor)


def require_text_fit(text: str, width_emu, height_emu, font_size: float, label: str, margin_in: float = 0.02):
    width_in = max(0.1, float(width_emu) / 914400.0 - margin_in * 2)
    height_in = max(0.01, float(height_emu) / 914400.0 - margin_in * 2)
    lines = estimate_wrapped_lines(text, width_in, font_size)
    needed = lines * font_size / 72.0 * 1.20
    if needed > height_in * 1.05:
        raise ValueError(
            f"{label} needs about {needed:.2f}in but the text box has {height_in:.2f}in; "
            "shorten the copy, split the slide, or increase the layout height"
        )


def fill_bg(slide, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_hex)


def add_rect(slide, x, y, w, h, fill, line=None, width=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    font_size: float,
    color_hex: str,
    *,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.02,
    line_spacing: float | None = None,
):
    font_size = safe_font_size(font_size)
    require_text_fit(text, w, h, font_size, f"text '{text[:24]}'", margin_in=margin)
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run_font = p.font
    run_font.name = FONT_CJK if has_cjk(text) else FONT
    run_font.size = Pt(font_size)
    run_font.bold = bold
    run_font.color.rgb = rgb(color_hex)
    if p.runs:
        apply_typeface(p.runs[0], text)
    return box


def has_cjk(text: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in text)


def apply_typeface(run, text: str) -> None:
    face = FONT_CJK if has_cjk(text) else FONT
    run.font.name = face
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        child = r_pr.find(qn(tag))
        if child is None:
            child = OxmlElement(tag)
            r_pr.append(child)
        child.set("typeface", face)


def safe_label(text: str) -> str:
    first = re.split(r"[：:、,，\s]", text.strip(), maxsplit=1)[0]
    return first[:8] if has_cjk(first) else first[:14].upper()


def density_units(text: str) -> float:
    units = 0.0
    for ch in str(text):
        if "\u4e00" <= ch <= "\u9fff":
            units += 1.0
        elif ch.isspace():
            units += 0.33
        elif ord(ch) < 128:
            units += 0.55
        else:
            units += 0.80
    return units


def validate_density(slides: list[SlideSpec], *, allow_dense: bool = False) -> None:
    if allow_dense:
        return
    problems: list[str] = []
    for idx, spec in enumerate(slides, 1):
        if idx == 1:
            continue
        bullets = spec.bullets or spec.body[1:]
        if density_units(spec.title) > 46:
            problems.append(f"slide {idx}: title is too long for the fixed header")
        if spec.body and density_units(spec.body[0]) > 72:
            problems.append(f"slide {idx}: takeaway line is too long; shorten it or let the visual carry the point")
        if len(bullets) > 5:
            problems.append(f"slide {idx}: {len(bullets)} bullets is too dense; split into another slide")
        for b_idx, bullet in enumerate(bullets, 1):
            if density_units(bullet) > 72:
                problems.append(f"slide {idx} bullet {b_idx}: too long for customer-facing layout")
    if problems:
        joined = "\n  - ".join(problems)
        raise SystemExit("Markdown is too dense for safe PPTX conversion:\n  - " + joined + "\nUse --allow-dense only for internal drafts.")


def slide_kicker(title: str) -> str:
    cjk = has_cjk(title)
    if any(key in title for key in ["下一", "步驟", "流程", "時程", "排程"]):
        return "下一步" if cjk else "ACTION PLAN"
    if any(key in title for key in ["比較", "選型", "方案", "差異"]):
        return "比較" if cjk else "DECISION"
    if any(key in title for key in ["儲位", "動線", "配置", "VNA", "倉儲"]):
        return "現場" if cjk else "LAYOUT IMPACT"
    if any(key in title for key in ["規格", "荷重", "地坪", "尺寸"]):
        return "確認" if cjk else "SPEC CHECK"
    return "重點" if cjk else "PROJECT NOTE"


def add_header(slide, spec: SlideSpec, theme: dict[str, str], idx: int) -> None:
    add_textbox(
        slide,
        slide_kicker(spec.title),
        Inches(0.72),
        Inches(0.42),
        Inches(2.05),
        Inches(0.26),
        8.5,
        theme["accent"],
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    add_textbox(slide, spec.title, Inches(0.72), Inches(0.72), Inches(10.15), Inches(0.58), 25, theme["ink"], bold=True)
    add_rect(slide, Inches(0.72), Inches(1.38), Inches(11.88), Inches(0.018), theme["rule"])
    add_textbox(
        slide,
        f"{idx:02d}",
        Inches(12.02),
        Inches(0.58),
        Inches(0.5),
        Inches(0.25),
        9,
        theme["muted"],
        bold=True,
        align=PP_ALIGN.RIGHT,
        margin=0,
    )


def add_cover(slide, spec: SlideSpec, deck_title: str, theme: dict[str, str], footer: str) -> None:
    fill_bg(slide, theme["deep"])
    add_rect(slide, Inches(0), Inches(0), Inches(0.28), SLIDE_H, theme["accent"])
    add_rect(slide, Inches(0.28), Inches(0), Inches(0.05), SLIDE_H, theme["accent2"])
    add_textbox(slide, footer, Inches(0.85), Inches(0.66), Inches(4.2), Inches(0.24), 9, "C9D5CF", bold=True, margin=0)
    valid_images = [path for path in spec.images if path.exists()]
    title_width = Inches(7.05) if valid_images else Inches(9.8)
    add_textbox(slide, deck_title or spec.title, Inches(0.82), Inches(1.75), title_width, Inches(1.45), 43, "FFFFFF", bold=True, line_spacing=0.9)
    # Keep the cover to one short line; supporting detail belongs in the deck,
    # not in a paragraph under the title.
    subtitle = spec.body[0] if spec.body else ""
    if subtitle:
        subtitle_width = Inches(6.55) if valid_images else Inches(7.9)
        add_textbox(slide, subtitle, Inches(0.88), Inches(3.32), subtitle_width, Inches(0.62), 17, "E5ECE8", margin=0.03)

    if valid_images:
        image_x, image_y = Inches(8.35), Inches(1.18)
        image_w, image_h = Inches(4.25), Inches(4.95)
        picture = slide.shapes.add_picture(str(valid_images[0]), image_x, image_y, width=image_w)
        if picture.height > image_h:
            old_width, old_height = picture.width, picture.height
            picture.height = image_h
            picture.width = int(old_width * image_h / old_height)
        picture.left = image_x + int((image_w - picture.width) / 2)
        picture.top = image_y + int((image_h - picture.height) / 2)

    add_rect(slide, Inches(0.86), Inches(5.28), Inches(3.4), Inches(0.64), "213F38", "43675C")
    add_textbox(slide, footer, Inches(1.08), Inches(5.47), Inches(3.0), Inches(0.2), 10, "DCE7E2", bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)


def add_takeaway_strip(slide, spec: SlideSpec, theme: dict[str, str]) -> None:
    # Keep the takeaway strip optional and short; the visual carries the detail.
    claim = spec.body[0] if spec.body else ""
    if not claim.strip():
        return
    add_rect(slide, Inches(0.72), Inches(1.88), Inches(12.0), Inches(0.70), theme["panel2"], None)
    add_rect(slide, Inches(0.72), Inches(1.88), Inches(0.09), Inches(0.70), theme["accent"])
    add_textbox(slide, claim, Inches(0.96), Inches(2.05), Inches(10.85), Inches(0.36), 16, theme["ink"], bold=True, valign=MSO_ANCHOR.MIDDLE)


def add_bullet_card(slide, item: str, index: int, x, y, w, h, theme: dict[str, str]) -> None:
    add_rect(slide, x, y, w, h, theme["panel"], theme["rule"], width=0.6)
    add_rect(slide, x, y, Inches(0.12), h, theme["accent"] if index == 1 else theme["accent2"])
    add_textbox(
        slide,
        f"{index:02d}",
        x + Inches(0.26),
        y + Inches(0.22),
        Inches(0.52),
        Inches(0.22),
        9,
        theme["accent"] if index == 1 else theme["accent2"],
        bold=True,
        margin=0,
    )
    label, detail = split_label_detail(item)
    add_textbox(slide, label, x + Inches(0.26), y + Inches(0.58), w - Inches(0.52), Inches(0.32), 15, theme["ink"], bold=True)
    add_textbox(slide, detail, x + Inches(0.26), y + Inches(1.02), w - Inches(0.58), h - Inches(1.18), 12.5, theme["text"])


def split_label_detail(item: str) -> tuple[str, str]:
    match = re.match(r"^([^：:]{2,18})[：:]\s*(.+)$", item)
    if match:
        return match.group(1).strip(), match.group(2).strip()
    if "，" in item:
        head, tail = item.split("，", 1)
        return head.strip(), tail.strip()
    if "," in item:
        head, tail = item.split(",", 1)
        return head.strip(), tail.strip()
    words = item.split()
    if len(words) > 9 and not has_cjk(item):
        return " ".join(words[:5]), " ".join(words[5:])
    return safe_label(item), item


def add_cards_layout(slide, spec: SlideSpec, theme: dict[str, str], idx: int) -> None:
    fill_bg(slide, theme["bg"])
    add_header(slide, spec, theme, idx)
    add_takeaway_strip(slide, spec, theme)

    bullets = spec.bullets or spec.body[1:]
    if not bullets:
        bullets = [spec.title]
    card_count = min(len(bullets), 4)
    y = Inches(3.08)
    if card_count <= 3:
        w = Inches(3.72)
        gap = Inches(0.28)
        for i, item in enumerate(bullets[:3], 1):
            add_bullet_card(slide, item, i, Inches(0.72) + (i - 1) * (w + gap), y, w, Inches(2.65), theme)
    else:
        w = Inches(5.85)
        h = Inches(1.72)
        positions = [
            (Inches(0.72), Inches(3.0)),
            (Inches(6.88), Inches(3.0)),
            (Inches(0.72), Inches(4.76)),
            (Inches(6.88), Inches(4.76)),
        ]
        for i, item in enumerate(bullets[:4], 1):
            add_bullet_card(slide, item, i, positions[i - 1][0], positions[i - 1][1], w, h, theme)

    add_footer(slide, theme)


def add_steps_layout(slide, spec: SlideSpec, theme: dict[str, str], idx: int) -> None:
    fill_bg(slide, theme["bg"])
    add_header(slide, spec, theme, idx)
    add_takeaway_strip(slide, spec, theme)

    bullets = spec.bullets or spec.body[1:] or [spec.title]
    x0 = Inches(0.78)
    y0 = Inches(3.05)
    step_w = Inches(3.62)
    step_h = Inches(2.35)
    gap = Inches(0.38)
    for i, item in enumerate(bullets[:3], 1):
        x = x0 + (i - 1) * (step_w + gap)
        add_rect(slide, x, y0, step_w, step_h, theme["panel"], theme["rule"], width=0.6)
        add_rect(slide, x, y0, step_w, Inches(0.58), theme["ink"])
        step_label = f"第 {i} 步" if has_cjk(spec.title) else f"STEP {i}"
        add_textbox(slide, step_label, x + Inches(0.22), y0 + Inches(0.17), Inches(1.4), Inches(0.18), 9.2, "FFFFFF", bold=True, margin=0)
        label, detail = split_label_detail(item)
        add_textbox(slide, label, x + Inches(0.28), y0 + Inches(0.82), step_w - Inches(0.56), Inches(0.36), 15, theme["ink"], bold=True)
        add_textbox(slide, detail, x + Inches(0.28), y0 + Inches(1.28), step_w - Inches(0.56), Inches(0.78), 12.5, theme["text"])
        if i < min(len(bullets), 3):
            add_textbox(slide, ">", x + step_w + Inches(0.12), y0 + Inches(1.02), Inches(0.18), Inches(0.22), 16, theme["accent"], bold=True, margin=0)

    if len(bullets) > 3:
        remainder = "\n".join(f"• {item}" for item in bullets[3:5])
        label = "補充確認項目：" if has_cjk(remainder) else "Additional checks: "
        add_textbox(slide, f"{label}\n{remainder}", Inches(0.86), Inches(5.82), Inches(10.9), Inches(0.72), 11, theme["muted"])
    add_footer(slide, theme)


def add_image_layout(slide, spec: SlideSpec, theme: dict[str, str], idx: int) -> None:
    fill_bg(slide, theme["bg"])
    add_header(slide, spec, theme, idx)
    add_takeaway_strip(slide, spec, theme)
    valid_images = [p for p in spec.images if p.exists()]
    if valid_images:
        image_x, image_y = Inches(7.92), Inches(3.0)
        image_w, image_h = Inches(4.55), Inches(3.55)
        pic = slide.shapes.add_picture(str(valid_images[0]), image_x, image_y, width=image_w)
        # Portrait images keep their aspect ratio; cap height so they stay
        # above the footer.
        if pic.height > image_h:
            old_w, old_h = pic.width, pic.height
            pic.height = image_h
            pic.width = int(old_w * image_h / old_h)
        pic.left = image_x + int((image_w - pic.width) / 2)
        pic.top = image_y + int((image_h - pic.height) / 2)
    bullets = spec.bullets or spec.body[1:]
    for i, item in enumerate(bullets[:4], 1):
        add_bullet_row(slide, item, i, Inches(0.88), Inches(3.05) + Inches(0.78) * (i - 1), Inches(6.5), theme)
    add_footer(slide, theme)


def add_bullet_row(slide, item: str, index: int, x, y, w, theme: dict[str, str], text_h: float = 0.48) -> None:
    add_rect(slide, x, y, Inches(0.34), Inches(0.34), theme["accent"] if index == 1 else theme["accent2"])
    add_textbox(slide, str(index), x + Inches(0.08), y + Inches(0.075), Inches(0.16), Inches(0.12), 8.5, "FFFFFF", bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_textbox(slide, item, x + Inches(0.52), y - Inches(0.03), w - Inches(0.5), Inches(text_h), 13.3, theme["text"])


def add_list_layout(slide, spec: SlideSpec, theme: dict[str, str], idx: int) -> None:
    fill_bg(slide, theme["bg"])
    add_header(slide, spec, theme, idx)
    add_takeaway_strip(slide, spec, theme)
    bullets = (spec.bullets or spec.body[1:] or [spec.title])[:5]
    # Row pitch must exceed the row text height so adjacent rows never overlap,
    # and the panel must enclose the last row (qa_check flags both).
    pitch = 0.58
    text_h = 0.50
    panel_h = 0.34 + pitch * (len(bullets) - 1) + text_h + 0.18
    add_rect(slide, Inches(0.82), Inches(3.0), Inches(11.7), Inches(panel_h), theme["panel"], theme["rule"], width=0.6)
    for i, item in enumerate(bullets, 1):
        add_bullet_row(slide, item, i, Inches(1.12), Inches(3.34) + Inches(pitch) * (i - 1), Inches(10.6), theme, text_h=text_h)
    add_footer(slide, theme)


def add_footer(slide, theme: dict[str, str]) -> None:
    add_textbox(slide, theme.get("footer", "TWS 奔騰物流"), Inches(0.72), Inches(6.86), Inches(4.1), Inches(0.18), 8.5, theme["muted"], margin=0)
    add_rect(slide, Inches(11.66), Inches(6.88), Inches(0.92), Inches(0.06), theme["accent"])


def content_kind(spec: SlideSpec) -> str:
    if any(p.exists() for p in spec.images):
        return "image"
    title = spec.title
    if any(key in title for key in ["下一", "步驟", "流程", "時程", "排程", "動線", "串接", "回報", "流向", "作業"]):
        return "steps"
    if 1 <= len(spec.bullets) <= 4:
        return "cards"
    return "list"


def build_deck(slides: list[SlideSpec], output: Path, title: str, theme_name: str, footer: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    theme = dict(THEMES[theme_name], footer=footer)

    for idx, spec in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        if idx == 1:
            add_cover(slide, spec, title or spec.title, theme, footer)
            continue
        kind = content_kind(spec)
        if kind == "steps":
            add_steps_layout(slide, spec, theme, idx)
        elif kind == "image":
            add_image_layout(slide, spec, theme, idx)
        elif kind == "cards":
            add_cards_layout(slide, spec, theme, idx)
        else:
            add_list_layout(slide, spec, theme, idx)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("input_md", type=Path)
    parser.add_argument("output_pptx", type=Path)
    parser.add_argument("--title", default="")
    parser.add_argument("--theme", choices=sorted(THEMES), default="industrial")
    parser.add_argument("--footer", default="TWS 奔騰物流",
                        help="brand/footer label; keep it customer-facing (no draft/review wording)")
    parser.add_argument("--allow-dense", action="store_true",
                        help="allow over-dense slides; only for internal drafts")
    args = parser.parse_args()

    markdown = args.input_md.read_text(encoding="utf-8")
    specs = [parse_slide(chunk, args.input_md.parent) for chunk in split_slides(markdown)]
    if not specs:
        raise SystemExit("No slides found in Markdown input.")
    validate_density(specs, allow_dense=args.allow_dense)
    build_deck(specs, args.output_pptx, args.title, args.theme, args.footer)
    print(f"Wrote {args.output_pptx} ({len(specs)} slides)")


if __name__ == "__main__":
    main()
