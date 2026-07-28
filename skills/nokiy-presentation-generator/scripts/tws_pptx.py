# -*- coding: utf-8 -*-
"""Shared python-pptx layout library for TWS customer decks.

Import this from deck build scripts instead of re-implementing helpers:

    import sys
    sys.path.insert(0, "<skill-root>/scripts")
    from tws_pptx import *

Design system: quiet industrial business style — off-white page, TWS dark
green + charcoal text, one warm orange accent, white cards, hairline borders.

Hard rules encoded here (do not bypass):
  - East-Asian and Latin typefaces are both set on every run.
  - Font sizes are clamped to readable minimums (body >= 10.5pt, caption >= 8.5pt)
    with a stderr warning when a caller asks for less.
  - All shapes have PowerPoint default shadows disabled.
  - Content must start below TITLE_ZONE (1.85in); add_slide() returns that y.
"""

from __future__ import annotations

import math
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---- canvas ----
PAGE_W_IN, PAGE_H_IN = 13.333333, 7.5
TITLE_ZONE = 1.85          # protected header depth (in); body content starts here
FOOTER_Y = 7.06
MARGIN = 0.68

# ---- color tokens ----
BG          = RGBColor(247, 246, 239)
GREEN       = RGBColor(15, 73, 58)     # titles
GREEN_2     = RGBColor(31, 98, 78)
GREEN_PALE  = RGBColor(231, 242, 235)
ORANGE      = RGBColor(214, 113, 47)   # single warm accent
ORANGE_PALE = RGBColor(248, 231, 217)
TEXT        = RGBColor(42, 55, 49)
MUTED       = RGBColor(108, 119, 113)
BORDER      = RGBColor(214, 222, 216)
WHITE       = RGBColor(255, 255, 255)
LIGHT       = RGBColor(250, 250, 246)

FONT_LATIN = "Helvetica Neue"
FONT_CJK = "PingFang TC"

# ---- readable minimums ----
MIN_BODY = 10.5
MIN_CAPTION = 8.5


def _clamp(size: float, floor: float = MIN_CAPTION) -> float:
    if size < floor:
        print(f"[tws_pptx] font {size}pt below minimum, clamped to {floor}pt", file=sys.stderr)
        return floor
    return size


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(PAGE_W_IN)
    prs.slide_height = Inches(PAGE_H_IN)
    return prs


def _set_run_fonts(run):
    """Set both Latin and East-Asian typefaces (font.name alone misses a:ea)."""
    rPr = run._r.get_or_add_rPr()
    for tag, face in (("a:latin", FONT_LATIN), ("a:ea", FONT_CJK), ("a:cs", FONT_LATIN)):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", face)


def add_text(slide, x, y, w, h, text, size, color=TEXT, bold=False,
             align=None, valign=None, line_spacing=None, floor=None):
    """Add a textbox. `text` may be a str or a list of line specs; each spec is
    a str or a tuple (str[, size[, color[, bold]]]) — omitted fields default."""
    if floor is None:
        floor = MIN_CAPTION if size < MIN_BODY else MIN_BODY
    fit_text = "\n".join(str(spec[0] if isinstance(spec, tuple) else spec) for spec in text) if isinstance(text, list) else str(text)
    require_text_fit(fit_text, w, h, max(size, floor), f"text '{fit_text[:24]}'", 1.0)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if valign:
        tf.vertical_anchor = valign
    lines = text if isinstance(text, list) else [(text, size, color, bold)]
    defaults = (size, color, bold)
    for i, spec in enumerate(lines):
        if isinstance(spec, tuple):
            t, s, c, b = (spec + defaults[len(spec) - 1:])[:4]
        else:
            t, s, c, b = spec, size, color, bold
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align:
            p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = t
        r.font.size = Pt(_clamp(s, floor))
        r.font.bold = b
        r.font.color.rgb = c
        _set_run_fonts(r)
    return box


def _visual_units(text: str) -> float:
    """Approximate rendered line length; CJK glyphs consume more width."""
    units = 0.0
    for ch in str(text):
        if "\u4e00" <= ch <= "\u9fff":
            units += 1.0
        elif ch.isspace():
            units += 0.33
        elif ord(ch) < 128:
            units += 0.55
        else:
            units += 0.80
    return units


def estimate_wrapped_lines(text: str, width_in: float, font_pt: float) -> int:
    """Conservative line estimate for text boxes before PowerPoint renders."""
    capacity = max(4.0, width_in * 72.0 / max(1.0, font_pt) * 0.92)
    lines = 0
    for raw in str(text).splitlines() or [""]:
        units = max(1.0, _visual_units(raw))
        lines += max(1, math.ceil(units / capacity))
    return lines


def required_text_height(text: str, width_in: float, font_pt: float, line_spacing: float = 1.20) -> float:
    """Estimate the vertical space needed by editable text in inches."""
    effective_size = max(1.0, font_pt)
    lines = sum(
        estimate_wrapped_lines(line, width_in, effective_size)
        for line in str(text).splitlines() or [""]
    )
    return max(1, lines) * effective_size / 72.0 * line_spacing


def require_text_fit(text: str, width_in: float, height_in: float, font_pt: float,
                     label: str, line_spacing: float = 1.20, tolerance: float = 1.05):
    """Fail early when a fixed-height helper cannot contain its requested copy."""
    needed = required_text_height(text, width_in, font_pt, line_spacing)
    if needed > max(0.01, height_in) * tolerance:
        raise ValueError(
            f"{label} needs about {needed:.2f}in but the text box has {height_in:.2f}in; "
            "shorten the copy, split the slide, or increase the layout height"
        )


def add_rrect(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=True, line_w=0.75):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def set_fill_alpha(shape, pct: int):
    """Real fill transparency (pct = 0 opaque … 100 invisible).
    Note: FillFormat has no .transparency in python-pptx — assigning it is a
    silent no-op. Use this XML helper instead."""
    alpha = str(int((100 - pct) * 1000))
    sp = shape.fill._xPr.find(qn("a:solidFill"))
    if sp is None:
        return
    clr = sp[0]
    for old in clr.findall(qn("a:alpha")):
        clr.remove(old)
    a = clr.makeelement(qn("a:alpha"), {"val": alpha})
    clr.append(a)


def add_slide(prs, section, page, title, subtitle=None, total=None,
              footer="TWS 奔騰物流"):
    """Standard content slide chrome. Returns (slide, content_top_y)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_text(slide, MARGIN, 0.28, 7.6, 0.24, "TWS 奔騰物流｜" + section, 9, GREEN_2, True)
    page_label = f"{page:02d}" + (f" / {total:02d}" if total else "")
    add_text(slide, PAGE_W_IN - MARGIN - 1.3, 0.28, 1.3, 0.24, page_label, 9, ORANGE, True, PP_ALIGN.RIGHT)
    rule = add_rrect(slide, MARGIN, 0.62, PAGE_W_IN - 2 * MARGIN, 0.012, BORDER, None, radius=False)
    add_text(slide, MARGIN, 0.82, PAGE_W_IN - 2 * MARGIN, 0.52, title, 22, GREEN, True)
    if subtitle:
        add_text(slide, MARGIN, 1.40, PAGE_W_IN - 2 * MARGIN, 0.32, subtitle, 10.5, MUTED)
    add_text(slide, MARGIN, FOOTER_Y, 6.0, 0.2, footer, 8.5, MUTED)
    return slide, TITLE_ZONE


def card_required_height(body_lines, width_in=3.0, body_size=MIN_BODY, title_size=12.5, title=""):
    """Minimum card height in inches for the default eyebrow/title/body layout."""
    body_lines = list(body_lines)
    body_size = _clamp(body_size, MIN_BODY)
    title_size = _clamp(title_size, MIN_BODY)
    line_h = (body_size / 72.0) * 1.28
    title_lines = estimate_wrapped_lines(title, width_in, title_size)
    wrapped_body_lines = sum(estimate_wrapped_lines(line, width_in, body_size) for line in body_lines) or 1
    title_h = max(0.30, title_lines * (title_size / 72.0) * 1.20)
    return 0.16 + 0.18 + 0.06 + title_h + 0.14 + wrapped_body_lines * line_h + 0.22


def add_card(slide, x, y, w, h, eyebrow, title, body_lines, fill=WHITE,
             title_size=12.5, body_size=MIN_BODY):
    body_lines = list(body_lines)
    title_size = _clamp(title_size, MIN_BODY)
    body_size = _clamp(body_size, MIN_BODY)
    # Keep card text from hugging or crossing the bottom edge. This expands the
    # card when callers underspecify height for the requested copy.
    min_h = card_required_height(body_lines, max(0.5, w - 0.4), body_size, title_size, title)
    if h < min_h:
        print(f"[tws_pptx] card height {h:.2f}in too small for {len(body_lines)} body line(s), expanded to {min_h:.2f}in", file=sys.stderr)
        h = min_h
    add_rrect(slide, x, y, w, h, fill, BORDER)
    add_text(slide, x + 0.2, y + 0.16, w - 0.4, 0.18, eyebrow, 8.5, ORANGE, True)
    title_h = max(0.30, estimate_wrapped_lines(title, max(0.5, w - 0.4), title_size) * (title_size / 72.0) * 1.20)
    body_y = y + 0.40 + title_h + 0.14
    body_h = h - (0.40 + title_h + 0.14 + 0.22)
    add_text(slide, x + 0.2, y + 0.40, w - 0.4, title_h, title, title_size, GREEN, True, floor=MIN_BODY)
    add_text(slide, x + 0.2, body_y, w - 0.4, body_h,
             [(l, body_size, TEXT, False) for l in body_lines], body_size,
             line_spacing=1.25, floor=MIN_BODY)
    return h


def add_card_stack(slide, x, y, w, cards, gap=0.30, default_h=1.20,
                   max_y=FOOTER_Y - 0.18, on_overflow="warn"):
    """Add a vertical stack of cards using each card's actual expanded height.

    cards entries are `(eyebrow, title, body_lines[, fill[, h]])`. Use this
    instead of fixed y offsets whenever cards are vertically stacked.
    """
    yy = y
    for spec in cards:
        eyebrow, title, body_lines = spec[:3]
        fill = spec[3] if len(spec) >= 4 else WHITE
        h = spec[4] if len(spec) >= 5 else default_h
        actual_h = add_card(slide, x, yy, w, h, eyebrow, title, body_lines, fill=fill)
        bottom = yy + actual_h
        if max_y is not None and bottom > max_y:
            msg = f"[tws_pptx] card stack bottom {bottom:.2f}in exceeds safe bottom {max_y:.2f}in; split the slide or shorten copy"
            if on_overflow == "raise":
                raise ValueError(msg)
            print(msg, file=sys.stderr)
        yy += actual_h + gap
    return yy


def add_badge(slide, x, y, text, color=ORANGE, w=0.42, h=0.26):
    sh = add_rrect(slide, x, y, w, h, color, None)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(_clamp(8.5))
    r.font.bold = True
    r.font.color.rgb = WHITE
    _set_run_fonts(r)
    return sh


def add_arrow(slide, x, y, w=0.32, h=0.13, color=ORANGE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_step(slide, x, y, w, h, no, title, body, fill=WHITE):
    require_text_fit(title, w - 0.9, 0.26, 11.5, f"step {no} title")
    require_text_fit(body, w - 0.4, h - 0.7, MIN_BODY, f"step {no} body", 1.25)
    add_rrect(slide, x, y, w, h, fill, BORDER)
    add_badge(slide, x + 0.18, y + 0.18, no)
    add_text(slide, x + 0.70, y + 0.18, w - 0.9, 0.26, title, 11.5, GREEN, True)
    add_text(slide, x + 0.2, y + 0.58, w - 0.4, h - 0.7, body, MIN_BODY, TEXT, line_spacing=1.25)


def add_flow(slide, items, x, y, w, h, gap=0.18):
    """items: [(title, body)] — horizontal flow with arrows."""
    n = len(items)
    box_w = (w - gap * (n - 1)) / n
    for i, (title, body) in enumerate(items):
        xx = x + i * (box_w + gap)
        require_text_fit(title, box_w - 0.24, 0.26, 11, f"flow item {i + 1} title")
        body_text = body if isinstance(body, str) else "\n".join(body)
        require_text_fit(body_text, box_w - 0.28, h - 0.68, MIN_BODY, f"flow item {i + 1} body", 1.20)
        add_rrect(slide, xx, y, box_w, h, WHITE, BORDER)
        add_text(slide, xx + 0.12, y + 0.18, box_w - 0.24, 0.26, title, 11, GREEN, True, PP_ALIGN.CENTER)
        add_text(slide, xx + 0.14, y + 0.55, box_w - 0.28, h - 0.68,
                 body.split("\n") if isinstance(body, str) else body,
                 MIN_BODY, TEXT, align=PP_ALIGN.CENTER, line_spacing=1.2)
        if i < n - 1:
            add_arrow(slide, xx + box_w + 0.02, y + h / 2 - 0.065, gap - 0.04, 0.13)


def add_banner(slide, x, y, w, text, fill=GREEN, color=WHITE, h=0.6, size=10.5):
    require_text_fit(text, w - 0.5, 0.24, size, "banner text")
    add_rrect(slide, x, y, w, h, fill, None)
    add_text(slide, x + 0.25, y + (h - 0.24) / 2, w - 0.5, 0.24, text, size, color, True, PP_ALIGN.CENTER)


def add_image(slide, path, x, y, w=None, h=None):
    if w is not None and h is not None:
        raise ValueError(
            "add_image cannot receive both width and height because that can distort the asset; "
            "use add_image_contain() or add_logo()"
        )
    kw = {}
    if w is not None:
        kw["width"] = Inches(w)
    if h is not None:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kw)


def add_image_contain(slide, path, x, y, w, h):
    """Place an image inside a box without cropping or distorting it."""
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    max_h = Inches(h)
    if pic.height > max_h:
        old_w, old_h = pic.width, pic.height
        pic.height = max_h
        pic.width = int(old_w * max_h / old_h)
    pic.left = Inches(x) + int((Inches(w) - pic.width) / 2)
    pic.top = Inches(y) + int((Inches(h) - pic.height) / 2)
    return pic


def add_logo(slide, path, x, y, w, h):
    """Logo-safe image placement: contain, no crop, no distortion."""
    return add_image_contain(slide, path, x, y, w, h)
