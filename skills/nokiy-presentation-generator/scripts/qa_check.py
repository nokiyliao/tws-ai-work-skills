#!/usr/bin/env python3
"""Mechanical QA gate for TWS customer decks.

Checks (in one pass over the PPTX):
  1. Package integrity (zip test).
  2. Slide count vs --expect-slides.
  3. Editable text exists (not an image-only deck).
  4. Banned / discouraged term scan from data/banned_terms.json (per --mode).
  5. Human-copy checks from data/copy_rules.json: title length, weak title
     patterns, and AI/marketing-fluff terms.
  6. Minimum explicit font size (default 10pt for customer-facing decks).
  7. Title-zone intrusion report: body shapes crossing or sitting inside the
     protected title/subtitle boundary (default 1.85in; cover skipped by default).
  8. Text/container containment and visible-object overlap reports.
  9. Locked-asset hash verification against ppt/media/*.
 10. Asset-registry provenance, generated-image role, and media-coverage checks.
 11. Optional asset-library selection manifest digest and role-boundary gate.
 12. Optional rendered-slide OCR scan for prompt/debug/watermark leakage.
 13. Embedded video report.

Exit code 0 = pass (warnings allowed), 1 = at least one failure.

Usage:
  uv run --with python-pptx python qa_check.py deck.pptx \
      --mode customer_facing --expect-slides 12 \
      --locked-asset assets/agv.png --skip-title-zone 1

Customer-facing mode treats title-zone intrusion, text fit, containment, visible-object overlap,
discouraged terms, and fonts below the minimum as failures by default.
Human-copy smell warnings are also failures in customer-facing mode.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
import subprocess
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from asset_snapshot import verify_materialized_selection

SKILL_ROOT = Path(__file__).resolve().parent.parent
DEFAULT_TERMS = SKILL_ROOT / "data" / "banned_terms.json"
DEFAULT_COPY_RULES = SKILL_ROOT / "data" / "copy_rules.json"
DEFAULT_OCR_HELPER = SKILL_ROOT / "scripts" / "ocr_rendered_slides.py"

VIDEO_EXT = {".mp4", ".mov", ".m4v", ".avi", ".wmv", ".mkv", ".webm"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".tif", ".tiff", ".bmp", ".webp", ".svg", ".emf", ".wmf", ".ico"}
REGISTRY_SOURCE_TYPES = {"user", "official", "local_authority", "generated", "unknown"}
SHA1_RE = re.compile(r"^[0-9a-f]{40}$", re.IGNORECASE)
GENERATED_FORBIDDEN_ROLES = {
    "logo",
    "brand",
    "product",
    "product_proof",
    "equipment_selection",
    "model_exterior",
    "real_site_evidence",
}

FAIL = "FAIL"
WARN = "WARN"
OK = "OK"


def is_ascii(term: str) -> bool:
    return all(ord(c) < 128 for c in term)


def term_hits(text: str, term: str) -> bool:
    if is_ascii(term):
        return re.search(r"(?i)(?<![A-Za-z0-9_])" + re.escape(term) + r"(?![A-Za-z0-9_])", text) is not None
    return term in text


def visual_units(text: str) -> float:
    units = 0.0
    for ch in str(text):
        if "\u4e00" <= ch <= "\u9fff":
            units += 1.0
        elif ch.isspace():
            units += 0.33
        elif ord(ch) < 128:
            units += 0.55
        else:
            units += 0.80
    return units


def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        try:
            is_group = sh.shape_type == 6  # MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            yield from iter_shapes(sh.shapes)


def shape_text(sh) -> str:
    if not getattr(sh, "has_text_frame", False):
        return ""
    return "\n".join(p.text for p in sh.text_frame.paragraphs)


def bbox_for(sh):
    try:
        top, left, w, h = sh.top, sh.left, sh.width, sh.height
    except Exception:
        return None
    if top is None or left is None or w is None or h is None or w <= 0 or h <= 0:
        return None
    return (int(left), int(top), int(left + w), int(top + h))


def bbox_area(box) -> int:
    return max(0, box[2] - box[0]) * max(0, box[3] - box[1])


def contains(outer, inner, pad=0) -> bool:
    return (
        outer[0] <= inner[0] + pad and
        outer[1] <= inner[1] + pad and
        outer[2] >= inner[2] - pad and
        outer[3] >= inner[3] - pad
    )


def soft_contains(outer, inner, slack=0) -> bool:
    """Return true when inner is inside outer, allowing a small visual bleed."""
    return (
        outer[0] - slack <= inner[0] and
        outer[1] - slack <= inner[1] and
        outer[2] + slack >= inner[2] and
        outer[3] + slack >= inner[3]
    )


def center_inside(outer, inner) -> bool:
    cx = (inner[0] + inner[2]) / 2
    cy = (inner[1] + inner[3]) / 2
    return outer[0] <= cx <= outer[2] and outer[1] <= cy <= outer[3]


def overlap_metrics(a, b):
    ix0 = max(a[0], b[0])
    iy0 = max(a[1], b[1])
    ix1 = min(a[2], b[2])
    iy1 = min(a[3], b[3])
    if ix1 <= ix0 or iy1 <= iy0:
        return 0, 0.0
    area = (ix1 - ix0) * (iy1 - iy0)
    smaller = max(1, min(bbox_area(a), bbox_area(b)))
    return area, area / smaller


def label_for(sh, txt: str) -> str:
    first = (txt.strip().splitlines() or [""])[0].strip()
    if first:
        return first[:24]
    name = getattr(sh, "name", "")
    return name or str(getattr(sh, "shape_type", "shape"))


def is_text_shape(sh, txt: str) -> bool:
    return bool(txt.strip()) and getattr(sh, "has_text_frame", False)


def _length_pt(value) -> float:
    try:
        return float(value.pt)
    except Exception:
        return 0.0


def text_fit_metrics(sh):
    """Estimate whether explicit text can fit its current textbox geometry."""
    if not getattr(sh, "has_text_frame", False):
        return None
    tf = sh.text_frame
    width_in = max(0.1, Emu(sh.width).inches)
    available_pt = max(0.0, Emu(sh.height).pt - _length_pt(tf.margin_top) - _length_pt(tf.margin_bottom))
    required_pt = 0.0
    saw_explicit_font = False
    total_lines = 0
    for p in tf.paragraphs:
        if not p.text.strip():
            continue
        sizes = []
        if p.font.size is not None:
            sizes.append(float(p.font.size.pt))
        for r in p.runs:
            if r.font.size is not None:
                sizes.append(float(r.font.size.pt))
        if not sizes:
            return None
        saw_explicit_font = True
        font_pt = max(sizes)
        capacity = max(4.0, width_in * 72.0 / max(1.0, font_pt) * 0.92)
        lines = max(1, int(math.ceil(visual_units(p.text) / capacity)))
        total_lines += lines
        required_pt += lines * font_pt * 1.08
    if not saw_explicit_font:
        return None
    return required_pt, available_pt, total_lines


def emu_area_from_square_inches(square_inches: float) -> int:
    return int(square_inches * 914400 * 914400)


def small_text_allowed(txt: str, box, zone_emu, footer_emu, min_caption: float, size_pt: float) -> bool:
    stripped = txt.strip()
    if not stripped:
        return True
    if size_pt < min_caption:
        return False
    if box is not None:
        top = box[1]
        if top < zone_emu or top >= footer_emu:
            return True
    return len(stripped) <= 8


def load_asset_registry(path: Path):
    data = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict) or not isinstance(data.get("assets"), list):
        raise ValueError("registry root must be an object with an 'assets' array")
    return data["assets"]


def validate_asset_registry(path: Path, registry_assets: list, image_media_hashes: dict[str, str]):
    """Validate registry shape, source policy, file provenance, and media coverage."""
    errors: list[str] = []
    registered_hashes: set[str] = set()
    seen_paths: set[str] = set()
    for i, item in enumerate(registry_assets, 1):
        if not isinstance(item, dict):
            errors.append(f"item {i} is not an object")
            continue
        missing = [key for key in ("deck_asset_path", "source_type", "role") if not str(item.get(key, "")).strip()]
        if missing:
            errors.append(f"item {i} missing required field(s): {', '.join(missing)}")
            continue
        asset_path = str(item["deck_asset_path"])
        source_type = str(item["source_type"]).strip().lower()
        role = str(item["role"]).strip().lower()
        if source_type not in REGISTRY_SOURCE_TYPES:
            errors.append(f"item {i} has invalid source_type '{source_type}'")
        if asset_path in seen_paths:
            errors.append(f"item {i} duplicates deck_asset_path '{asset_path}'")
        seen_paths.add(asset_path)
        allowed_use = normalize_allowed_use(item.get("allowed_use"))
        if source_type == "generated" and ({role} | allowed_use) & GENERATED_FORBIDDEN_ROLES:
            errors.append(f"generated asset '{asset_path}' cannot be used as {role or sorted(allowed_use)}")
        if source_type == "generated" and not str(item.get("prompt", "")).strip():
            errors.append(f"generated asset '{asset_path}' requires its generation prompt")
        p = Path(asset_path)
        if not p.is_absolute():
            p = path.parent / p
        try:
            actual_sha1 = hashlib.sha1(p.read_bytes()).hexdigest()
        except OSError as e:
            errors.append(f"cannot read {asset_path}: {e}")
            continue
        declared_sha1 = str(item.get("sha1", "")).strip().lower()
        if declared_sha1 and not SHA1_RE.fullmatch(declared_sha1):
            errors.append(f"invalid sha1 for {asset_path}")
        if declared_sha1 and declared_sha1 != actual_sha1:
            errors.append(f"sha1 mismatch for {asset_path} (registry {declared_sha1[:10]}…, file {actual_sha1[:10]}…)")
        if actual_sha1 not in image_media_hashes.values():
            errors.append(f"{asset_path} is not embedded in ppt/media")
        registered_hashes.add(actual_sha1)
    missing_media = set(image_media_hashes.values()) - registered_hashes
    if missing_media:
        errors.append(f"{len(missing_media)} embedded image asset(s) are not listed in the registry")
    return errors


def normalize_allowed_use(value) -> set[str]:
    if value is None:
        return set()
    if isinstance(value, str):
        return {value.lower()}
    return {str(v).lower() for v in value}


def run_ocr_helper(helper: Path, rendered_dir: Path) -> list[dict]:
    cmd = [sys.executable, str(helper), str(rendered_dir), "--json"]
    res = subprocess.run(cmd, text=True, capture_output=True)
    if res.returncode != 0:
        raise RuntimeError((res.stderr or res.stdout or "OCR helper failed").strip())
    return json.loads(res.stdout or "[]")


def max_font_pt(sh) -> float:
    max_pt = 0.0
    if not getattr(sh, "has_text_frame", False):
        return max_pt
    for p in sh.text_frame.paragraphs:
        p_size = p.font.size
        for r in p.runs:
            eff = r.font.size if r.font.size is not None else p_size
            if eff is not None:
                max_pt = max(max_pt, float(eff.pt))
    return max_pt


def load_copy_rules(path: Path, mode: str) -> dict:
    if not path.exists():
        return {}
    data = json.loads(path.read_text(encoding="utf-8"))
    return data.get(mode, data.get("customer_facing", {})) if isinstance(data, dict) else {}


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--mode", default="customer_facing")
    ap.add_argument("--expect-slides", type=int)
    ap.add_argument("--min-font", type=float, default=10.0, help="warn/fail when an explicit run size is below this (pt)")
    ap.add_argument("--title-zone", type=float, default=1.85, help="protected header depth in inches")
    ap.add_argument("--skip-title-zone", default="1", help="comma-separated slide numbers to skip zone check (default: 1 = cover)")
    ap.add_argument("--skip-overlap", default="", help="comma-separated slide numbers to skip overlap check")
    ap.add_argument("--layout-exception-reason", default="", help="required explanation when any non-default layout check is skipped")
    ap.add_argument("--overlap-threshold", type=float, default=0.04,
                    help="warn when overlap area exceeds this share of the smaller object")
    ap.add_argument("--overlap-min-area", type=float, default=0.03,
                    help="ignore overlaps smaller than this many square inches")
    ap.add_argument("--locked-asset", action="append", default=[], type=Path)
    ap.add_argument("--banned-terms", type=Path, default=DEFAULT_TERMS)
    ap.add_argument("--copy-rules", type=Path, default=DEFAULT_COPY_RULES)
    ap.add_argument("--ban", action="append", default=[], help="extra banned term (repeatable)")
    ap.add_argument("--allow-video", action="store_true")
    ap.add_argument("--asset-registry", type=Path,
                    help="JSON registry for deck assets with source_type, role, allowed_use, sha1, and deck_asset_path")
    ap.add_argument("--asset-selection-manifest", type=Path,
                    help="selection manifest to verify before QA uses library assets")
    ap.add_argument("--asset-library", type=Path,
                    help="job-local materialized asset snapshot containing catalog.json")
    ap.add_argument("--build-note", type=Path,
                    help="traceability note for the build: audience, authority inputs, assets, output, and QA level")
    ap.add_argument("--rendered-dir", type=Path, help="Folder containing rendered slide PNG/JPG images for OCR")
    ap.add_argument("--ocr-rendered", action="store_true", help="OCR rendered slides and scan for banned terms/prompt leakage")
    ap.add_argument("--ocr-helper", type=Path, default=DEFAULT_OCR_HELPER)
    ap.add_argument("--strict-zone", action="store_true")
    ap.add_argument("--strict-overlap", action="store_true")
    ap.add_argument("--strict-discouraged", action="store_true")
    ap.add_argument("--strict-copy", action="store_true")
    ap.add_argument("--strict-fit", action="store_true")
    args = ap.parse_args()

    if args.mode == "customer_facing":
        args.strict_zone = True
        args.strict_overlap = True
        args.strict_discouraged = True
        args.strict_copy = True
        args.strict_fit = True
        args.min_font = max(args.min_font, 10.0)

    results: list[tuple[str, str]] = []  # (level, message)

    def add(level: str, msg: str):
        results.append((level, msg))

    # Fail closed before opening the deck when a library selection is supplied.
    if bool(args.asset_selection_manifest) != bool(args.asset_library):
        print("FAIL asset-selection: --asset-selection-manifest and --asset-library must be provided together")
        return 1
    if args.asset_selection_manifest:
        try:
            detail = verify_materialized_selection(args.asset_library, args.asset_selection_manifest)
        except (OSError, ValueError, KeyError, json.JSONDecodeError) as exc:
            print(f"FAIL asset-selection: {exc}")
            return 1
        add(OK, f"asset-selection: {detail}")

    # 1. package integrity
    try:
        zf = zipfile.ZipFile(args.pptx)
        bad = zf.testzip()
        if bad:
            add(FAIL, f"package: corrupt member {bad}")
        else:
            add(OK, "package: zip integrity ok")
    except Exception as e:
        print(f"FAIL package: cannot open {args.pptx}: {e}")
        return 1

    media = [n for n in zf.namelist() if n.startswith("ppt/media/")]
    media_hashes = {}
    for n in media:
        media_hashes[n] = hashlib.sha1(zf.read(n)).hexdigest()

    image_media_hashes = {
        n: h for n, h in media_hashes.items()
        if Path(n).suffix.lower() in IMAGE_EXT
    }

    # traceability and asset registry
    if args.mode == "customer_facing":
        if not args.build_note:
            add(FAIL, "traceability: --build-note is required for customer-facing QA")
        elif not args.build_note.exists() or not args.build_note.read_text(encoding="utf-8").strip():
            add(FAIL, f"traceability: build note missing or empty: {args.build_note}")
        else:
            add(OK, f"traceability: build note present ({args.build_note})")
        if image_media_hashes and not args.asset_registry:
            add(FAIL, "asset-registry: --asset-registry is required when the deck contains images")
    elif args.build_note:
        if args.build_note.exists() and args.build_note.read_text(encoding="utf-8").strip():
            add(OK, f"traceability: build note present ({args.build_note})")
        else:
            add(WARN, f"traceability: build note missing or empty: {args.build_note}")

    # videos
    vids = [n for n in media if Path(n).suffix.lower() in VIDEO_EXT]
    if vids:
        add(WARN if args.allow_video else FAIL,
            f"media: embedded video found: {', '.join(vids)}" + ("" if args.allow_video else " (use --allow-video if intended)"))
    else:
        add(OK, "media: no embedded video")

    # locked assets
    for asset in args.locked_asset:
        try:
            h = hashlib.sha1(asset.read_bytes()).hexdigest()
        except OSError as e:
            add(FAIL, f"locked-asset: cannot read {asset}: {e}")
            continue
        if h in media_hashes.values():
            add(OK, f"locked-asset: {asset.name} present in ppt/media")
        else:
            add(FAIL, f"locked-asset: {asset.name} NOT found in ppt/media (sha1 {h[:10]}…)")

    # asset registry
    if args.asset_registry:
        try:
            registry_assets = load_asset_registry(args.asset_registry)
        except Exception as e:
            registry_assets = []
            add(FAIL, f"asset-registry: cannot read {args.asset_registry}: {e}")
        if not registry_assets:
            add(FAIL, f"asset-registry: no assets listed in {args.asset_registry}")
        for error in validate_asset_registry(args.asset_registry, registry_assets, image_media_hashes):
            add(FAIL, f"asset-registry: {error}")
        if registry_assets and not any(level == FAIL and msg.startswith("asset-registry:") for level, msg in results):
            add(OK, f"asset-registry: {len(registry_assets)} asset(s) validated and cover embedded images")

    prs = Presentation(args.pptx)

    # 2. slide count
    n_slides = len(prs.slides)
    if args.expect_slides is not None:
        if n_slides == args.expect_slides:
            add(OK, f"slides: count {n_slides} matches expectation")
        else:
            add(FAIL, f"slides: count {n_slides}, expected {args.expect_slides}")
    else:
        add(OK, f"slides: count {n_slides} (no expectation given)")

    # term lists
    terms = {"banned": list(args.ban), "discouraged": []}
    if args.banned_terms.exists():
        data = json.loads(args.banned_terms.read_text(encoding="utf-8"))
        mode = data.get("modes", {}).get(args.mode, {})
        terms["banned"] += mode.get("banned", [])
        terms["discouraged"] += mode.get("discouraged", [])
    else:
        add(WARN, f"terms: list file not found: {args.banned_terms}")

    copy_rules = load_copy_rules(args.copy_rules, args.mode)
    max_title_units = float(copy_rules.get("max_title_units", 34))
    max_claim_units = float(copy_rules.get("max_claim_units", 72))
    max_slide_body_units = float(copy_rules.get("max_slide_body_units", 120))
    banned_fluff_terms = list(copy_rules.get("banned_fluff_terms", []))
    discouraged_fluff_terms = list(copy_rules.get("discouraged_fluff_terms", []))
    weak_title_patterns = list(copy_rules.get("weak_title_patterns", []))

    skip_zone = {int(x) for x in str(args.skip_title_zone).split(",") if x.strip()}
    skip_overlap = {int(x) for x in str(args.skip_overlap).split(",") if x.strip()}
    non_default_zone_skips = skip_zone - {1}
    if non_default_zone_skips or skip_overlap:
        reason = args.layout_exception_reason.strip()
        if args.mode == "customer_facing" and len(reason) < 10:
            add(FAIL, "layout exceptions: customer-facing skips require --layout-exception-reason (10+ characters)")
        else:
            add(WARN, f"layout exceptions active: title-zone={sorted(non_default_zone_skips)} overlap={sorted(skip_overlap)}; reason: {reason or 'not supplied'}")
    zone_emu = Emu(int(args.title_zone * 914400))
    footer_emu = Emu(int(6.78 * 914400))
    overlap_min_area = emu_area_from_square_inches(args.overlap_min_area)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    any_text = False
    small_fonts: list[str] = []
    zone_hits: list[str] = []
    overlap_hits: list[str] = []
    offslide_hits: list[str] = []
    containment_hits: list[str] = []
    fit_hits: list[str] = []
    unresolved_font_hits: list[str] = []
    banned_hits: list[str] = []
    discouraged_hits: list[str] = []
    ocr_banned_hits: list[str] = []
    ocr_discouraged_hits: list[str] = []
    ocr_copy_hits: list[str] = []
    title_hits: list[str] = []
    copy_banned_hits: list[str] = []
    copy_discouraged_hits: list[str] = []
    weak_copy_pattern_hits: list[str] = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_text_parts = []
        slide_body_units = 0.0
        overlap_candidates = []
        for sh in iter_shapes(slide.shapes):
            txt = shape_text(sh)
            box = bbox_for(sh)
            if txt.strip():
                any_text = True
                slide_text_parts.append(txt)
                units = visual_units(txt.strip())
                top = box[1] if box is not None else None
                if box is not None and top >= zone_emu and box[3] < footer_emu:
                    slide_body_units += units
                in_title_band = top is not None and idx != 1 and int(0.55 * 914400) <= top < int(args.title_zone * 914400)
                if in_title_band and max_font_pt(sh) >= 16:
                    title_line = (txt.strip().splitlines() or [""])[0].strip()
                    title_units = visual_units(title_line)
                    if title_units > max_title_units:
                        title_hits.append(f"slide {idx}: title too long ({title_units:.0f} units) 「{title_line[:34]}」")
                    for pattern in weak_title_patterns:
                        if re.search(pattern, title_line):
                            title_hits.append(f"slide {idx}: weak title pattern /{pattern}/ 「{title_line[:34]}」")
                elif units > max_claim_units and max_font_pt(sh) >= 14:
                    copy_discouraged_hits.append(f"slide {idx}: long claim text ({units:.0f} units) 「{txt.strip()[:34]}」")
                metrics = text_fit_metrics(sh)
                if metrics:
                    required_pt, available_pt, total_lines = metrics
                    if required_pt > available_pt * 1.10:
                        fit_hits.append(
                            f"slide {idx}: '{label_for(sh, txt)}' needs about {required_pt:.1f}pt but has {available_pt:.1f}pt ({total_lines} line(s))"
                        )
                if max_font_pt(sh) == 0:
                    unresolved_font_hits.append(f"slide {idx}: '{label_for(sh, txt)}' has no explicit run font size")
            # font sizes (run-level overrides paragraph-level default)
            if getattr(sh, "has_text_frame", False):
                for p in sh.text_frame.paragraphs:
                    p_size = p.font.size
                    for r in p.runs:
                        eff = r.font.size if r.font.size is not None else p_size
                        if eff is not None and eff.pt < args.min_font and r.text.strip():
                            if small_text_allowed(r.text, box, zone_emu, footer_emu, 8.5, eff.pt):
                                continue
                            small_fonts.append(f"slide {idx}: {eff.pt:.1f}pt 「{r.text.strip()[:18]}」")
            if box is None:
                continue
            left, top, right, bottom = box
            w = right - left
            h = bottom - top
            # ignore full-bleed backgrounds and hairlines
            if w >= slide_w * 0.92 and h >= slide_h * 0.92:
                continue
            if bbox_area(box) < overlap_min_area:
                continue
            label = label_for(sh, txt)
            tolerance = int(0.02 * 914400)
            if left < -tolerance or top < -tolerance or right > slide_w + tolerance or bottom > slide_h + tolerance:
                offslide_hits.append(
                    f"slide {idx}: '{label}' bbox extends outside slide "
                    f"(left={Emu(left).inches:.2f}in top={Emu(top).inches:.2f}in "
                    f"right={Emu(right).inches:.2f}in bottom={Emu(bottom).inches:.2f}in)")
            overlap_candidates.append((box, sh, txt, label))
            # title zone
            if idx not in skip_zone:
                # violation: shape starts above the zone boundary and extends below it
                if top < zone_emu and bottom > zone_emu:
                    zone_hits.append(
                        f"slide {idx}: '{label}' bbox top={Emu(top).inches:.2f}in h={Emu(h).inches:.2f}in crosses {args.title_zone}in")
                elif int(0.55 * 914400) <= top and bottom <= zone_emu and h > int(0.05 * 914400):
                    # Header text itself is expected here; body panels/images are not.
                    header_text = is_text_shape(sh, txt) and top < int(1.55 * 914400)
                    if not header_text:
                        zone_hits.append(
                            f"slide {idx}: '{label}' is fully inside the protected header band ({Emu(top).inches:.2f}–{Emu(bottom).inches:.2f}in)")

        for a_i in range(len(overlap_candidates)):
            a_box, _a_sh, _a_txt, a_label = overlap_candidates[a_i]
            for b_i in range(a_i + 1, len(overlap_candidates)):
                b_box, _b_sh, _b_txt, b_label = overlap_candidates[b_i]
                # Nested elements are normal in card/image composition.
                # Allow a 0.02in tolerance for near-nested geometry.
                pad = int(0.02 * 914400)
                if contains(a_box, b_box, pad) or contains(b_box, a_box, pad):
                    continue
                # Treat text placed on a panel/image as a container relation,
                # but report text that bleeds outside the container separately.
                relation = None
                for text_box, text_sh, text_label, outer_box, outer_sh, outer_label in (
                    (a_box, _a_sh, a_label, b_box, _b_sh, b_label),
                    (b_box, _b_sh, b_label, a_box, _a_sh, a_label),
                ):
                    if not is_text_shape(text_sh, shape_text(text_sh)):
                        continue
                    if is_text_shape(outer_sh, shape_text(outer_sh)):
                        continue
                    if bbox_area(text_box) >= bbox_area(outer_box) * 0.45:
                        continue
                    slack = int(0.12 * 914400)
                    if center_inside(outer_box, text_box) and soft_contains(outer_box, text_box, slack):
                        relation = (text_box, text_label, outer_box, outer_label)
                        break
                if relation:
                    text_box, text_label, outer_box, outer_label = relation
                    containment_hits.append(
                        f"slide {idx}: text '{text_label}' extends outside container '{outer_label}'"
                    )
                    continue
                if idx in skip_overlap:
                    continue
                area, ratio = overlap_metrics(a_box, b_box)
                if area >= overlap_min_area and ratio >= args.overlap_threshold:
                    overlap_hits.append(
                        f"slide {idx}: '{a_label}' overlaps '{b_label}' "
                        f"({ratio:.0%} of smaller object)")

        full_text = "\n".join(slide_text_parts)
        notes_text = ""
        if getattr(slide, "has_notes_slide", False):
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            if notes_frame is not None:
                notes_text = notes_frame.text or ""
        copy_review_text = "\n".join(part for part in (full_text, notes_text) if part)
        if slide_body_units > max_slide_body_units:
            copy_discouraged_hits.append(
                f"slide {idx}: body copy density {slide_body_units:.0f} units exceeds {max_slide_body_units:.0f}; replace prose with a visual or split the slide"
            )
        for t in terms["banned"]:
            if term_hits(full_text, t):
                banned_hits.append(f"slide {idx}: banned 「{t}」")
        for t in terms["discouraged"]:
            if term_hits(full_text, t):
                discouraged_hits.append(f"slide {idx}: discouraged 「{t}」")
        for t in banned_fluff_terms:
            if term_hits(copy_review_text, t):
                location = "visible copy / notes" if term_hits(notes_text, t) else "visible copy"
                copy_banned_hits.append(f"slide {idx}: fluff in {location} 「{t}」")
        for t in discouraged_fluff_terms:
            if term_hits(full_text, t):
                copy_discouraged_hits.append(f"slide {idx}: copy smell 「{t}」")
        for pattern in copy_rules.get("weak_copy_patterns", []):
            if re.search(pattern, copy_review_text):
                location = "visible copy / notes" if re.search(pattern, notes_text) else "visible copy"
                weak_copy_pattern_hits.append(f"slide {idx}: weak copy pattern in {location} /{pattern}/")

    if args.ocr_rendered:
        if not args.rendered_dir:
            add(FAIL, "ocr: --ocr-rendered requires --rendered-dir")
        elif not args.rendered_dir.exists():
            add(FAIL, f"ocr: rendered dir not found: {args.rendered_dir}")
        elif not args.ocr_helper.exists():
            add(FAIL, f"ocr: helper not found: {args.ocr_helper}")
        else:
            try:
                ocr_rows = run_ocr_helper(args.ocr_helper, args.rendered_dir)
                add(OK, f"ocr: scanned {len(ocr_rows)} rendered image(s)")
                for row in ocr_rows:
                    path = row.get("path", "rendered image")
                    if row.get("error"):
                        add(WARN, f"ocr: {path}: {row['error']}")
                        continue
                    text = row.get("text", "")
                    for t in terms["banned"]:
                        if term_hits(text, t):
                            ocr_banned_hits.append(f"{Path(path).name}: banned 「{t}」")
                    for t in terms["discouraged"]:
                        if term_hits(text, t):
                            ocr_discouraged_hits.append(f"{Path(path).name}: discouraged 「{t}」")
                    for t in banned_fluff_terms + discouraged_fluff_terms:
                        if term_hits(text, t):
                            ocr_copy_hits.append(f"{Path(path).name}: copy smell 「{t}」")
                    for pattern in copy_rules.get("weak_copy_patterns", []):
                        if re.search(pattern, text):
                            ocr_copy_hits.append(f"{Path(path).name}: weak copy pattern /{pattern}/")
            except Exception as e:
                add(FAIL, f"ocr: {e}")

    # 3. editable text
    add(OK if any_text else FAIL,
        "text: editable text present" if any_text else "text: no editable text found (image-only deck?)")

    # 4. terms
    if banned_hits:
        for h in banned_hits:
            add(FAIL, f"terms: {h}")
    else:
        add(OK, f"terms: no banned terms for mode '{args.mode}'")
    for h in discouraged_hits:
        add(FAIL if args.strict_discouraged else WARN, f"terms: {h}")
    for h in title_hits:
        add(FAIL if args.mode == "customer_facing" else WARN, f"copy: {h}")
    for h in copy_banned_hits:
        add(FAIL, f"copy: {h}")
    for h in copy_discouraged_hits[:30]:
        add(FAIL if args.strict_copy else WARN, f"copy: {h}")
    if len(copy_discouraged_hits) > 30:
        add(FAIL if args.strict_copy else WARN, f"copy: …and {len(copy_discouraged_hits) - 30} more copy smell warnings")
    for h in weak_copy_pattern_hits:
        add(FAIL if args.strict_copy else WARN, f"copy: {h}")
    for h in ocr_banned_hits:
        add(FAIL, f"ocr terms: {h}")
    for h in ocr_discouraged_hits:
        add(FAIL if args.strict_discouraged else WARN, f"ocr terms: {h}")
    for h in ocr_copy_hits:
        add(FAIL if args.strict_copy else WARN, f"ocr copy: {h}")

    # 5. fonts
    if small_fonts:
        for h in small_fonts[:20]:
            add(FAIL if args.mode == "customer_facing" else WARN, f"font: {h}")
        if len(small_fonts) > 20:
            add(FAIL if args.mode == "customer_facing" else WARN,
                f"font: …and {len(small_fonts) - 20} more runs below {args.min_font}pt")
    else:
        add(OK, f"font: no explicit run below {args.min_font}pt")
    for h in unresolved_font_hits[:20]:
        add(WARN, f"font: {h}")
    if len(unresolved_font_hits) > 20:
        add(WARN, f"font: …and {len(unresolved_font_hits) - 20} more text shape(s) without explicit font size")

    # 6. text fit
    if fit_hits:
        for h in fit_hits[:30]:
            add(FAIL if args.strict_fit else WARN, f"text-fit: {h}")
        if len(fit_hits) > 30:
            add(FAIL if args.strict_fit else WARN, f"text-fit: …and {len(fit_hits) - 30} more text boxes likely to clip or wrap")
    else:
        add(OK, "text-fit: no explicit text box is estimated to overflow")

    # 7. title zone
    if zone_hits:
        for h in zone_hits:
            add(FAIL if args.strict_zone else WARN, f"zone: {h}")
    else:
        add(OK, f"zone: no body shape intrudes into the {args.title_zone}in header boundary")

    # 8. container containment
    if containment_hits:
        for h in containment_hits[:30]:
            add(FAIL if args.strict_overlap else WARN, f"containment: {h}")
        if len(containment_hits) > 30:
            add(FAIL if args.strict_overlap else WARN, f"containment: …and {len(containment_hits) - 30} more text/container boundary issues")
    else:
        add(OK, "containment: no text extends beyond a likely panel/image container")

    # 9. visible-object overlap
    if overlap_hits:
        for h in overlap_hits[:30]:
            add(FAIL if args.strict_overlap else WARN, f"overlap: {h}")
        if len(overlap_hits) > 30:
            add(FAIL if args.strict_overlap else WARN, f"overlap: …and {len(overlap_hits) - 30} more geometry collisions")
    else:
        add(OK, "overlap: no visible-object geometry collisions above threshold")

    # 10. off-slide geometry
    if offslide_hits:
        for h in offslide_hits[:20]:
            add(FAIL if args.strict_overlap else WARN, f"off-slide: {h}")
        if len(offslide_hits) > 20:
            add(FAIL if args.strict_overlap else WARN, f"off-slide: …and {len(offslide_hits) - 20} more objects outside the slide")
    else:
        add(OK, "off-slide: no visible object extends outside the slide")

    # report
    fails = [m for l, m in results if l == FAIL]
    warns = [m for l, m in results if l == WARN]
    for level, msg in results:
        print(f"[{level}] {msg}")
    print()
    print(f"SUMMARY: {len(fails)} failure(s), {len(warns)} warning(s), {n_slides} slide(s) — "
          + ("FAIL" if fails else "PASS"))
    return 1 if fails else 0


if __name__ == "__main__":
    sys.exit(main())
